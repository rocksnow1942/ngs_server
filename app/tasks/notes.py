from rq import get_current_job
from app import db
from app.models import NGSSampleGroup,Primers,Rounds,Sequence,KnownSequence,Task,SeqRound,Analysis
from app import create_app
from flask import current_app
import json
from itertools import islice
from app.utils.ngs_util import reverse_comp, file_blocks, create_folder_if_not_exist,lev_distance
from collections import Counter
import re
from app.utils.analysis import DataReader
import os
from pptx import Presentation


app = create_app()
app.app_context().push()




file='/Users/hui/Documents/Scripts/pptconverter/othertests.pptx'


import hashlib
hasher = hashlib.md5()
with open(file, 'rb') as afile:
    buf = afile.read()
    hasher.update(buf)
print(hasher.hexdigest())


6f16f623495e686c1ec3503652128362


3626e7c8d74cf0baab0498a100c95f65



pj='/Users/hui/Aptitude_Cloud/Aptitude Users/R&D/Projects'
pj='/Users/hui/Cloudstation/R&D/Projects'


def glob_pptx(path):
    result = []
    for root, dirs, files in os.walk(path):
        for file in files:
            if file.endswith(".pptx") and ('conflict' not in file.lower()):
                  result.append(os.path.join(root, file))
    return result


f = glob_pptx(pj)[3]


def parese_presentation(f):
    ppt=Presentation(f)
    project = f.replace(pj,'').split('/')[1]
    presentation= os.path.basename(f)[0:-5]
    revision = ppt.core_properties.revision
    slides=[]
    hasher = hashlib.md5()
    with open(f, 'rb') as afile:
        buf = afile.read()
        hasher.update(buf)
    md5=hasher.hexdigest()
    for page,slide in enumerate(ppt.slides):
        temp=[]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                temp.append((shape.text,shape.top))
        if not temp: temp.append(("No Title",0))

        temp.sort(key=lambda x:x[1])

        slides.append(dict(key='Slide'+str(page+1),title=temp[0][0],body="\n".join([i[0] for i in temp[1:]])))
    return {'project':project, 'ppt':presentation, 'md5':md5, 'slides': slides}

result=[]
for i in glob_pptx(pj):
    result.append(parese_presentation(i))


len(result)
result[50]





project

deck

revision

for s in slides:
    for k,item in s.items():
        print(k)
        print(item.replace('\n', ''))




res=glob_pptx(pj)
os.path.basename(res[0])[0:-5]

for f in glob_pptx(pj):





    ppt.slides[0]

    ppt.core_properties.created

    ppt.core_properties.revision

    ppt2.core_properties.revision


    s3=ppt.slides[46]

    dir(s3)
    tt=s3.shapes[1]


    tt.top
    tt.text
    ct=s3.shapes[2]
    ct.top
    ct
    tx = s3.shapes[0]
    tx.top
    tx.text
