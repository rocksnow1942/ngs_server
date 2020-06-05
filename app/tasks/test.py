from pptx import Presentation
from dateutil import parser
from datetime import datetime
import re

import requests



r=requests.get('http://ams/api/get_plojo_data',json={'keys':['ams1234']})

r.json()

flag = "note"
regx = re.compile(f'<(?P<pt>{flag})>(.+)</(?P=pt)>')

date = re.compile("(?P<y>20\d{2}|\d{2})\W*(?P<m>[0]\d|1[0-2]|[1-9])\W*(?P<d>[0-2]\d|3[01]|\d)(?:\s*\((?P<author>[a-zA-Z\s]+)\))?")

date = re.compile("(?P<y>20\d{2}|\d{2})\W*(?P<m>[0]\d|1[0-2]|[1-9]|\d)\W*(?P<d>[0-2]\d|3[01]|\d)")

s = "2020121"
date.findall(s)


date.search('20120103  ( ) ').groups()

parser.parse('11/01/2010')

pptx = re.compile(r'(\S*)(\s*)\.pptx')
pptx.search('aba.1. pptx').groups()

path = '/Users/hui/Desktop/2020-02-03. .pptx'
re.sub(pptx,'\g<1>.pptx',name,)
import os
os.rename('/Users/hui/Desktop/2020-02-03. .pptx',re.sub(pptx,'\g<1>.pptx',name,))
newpath = re.sub(r'(\S*)(\s*)\.pptx','\g<1>.pptx',path)
newpath


dateAuthor = re.compile(
    r"\A\s*(?P<y>20\d{2}|\d{2})\W*(?P<m>[0]\d|1[0-2]|[1-9])\W*(?P<d>[0-2]\d|3[01]|\d)(?:\s*\((?P<author>[a-zA-Z\s]+)\))?")

title = "2011-17 GC24-Dolphin1 Dolphin2 Binding Test 0 3 10 30 100"

dateAuthor.search(title) and dateAuthor.search(title).groups()



def remove_pptx_blank(path):
    newpath = re.sub(r'(\S*)(\s*)\.pptx', r'\g<1>.pptx', path)
    if path != newpath and os.path.exists(path):
        repl = r'\g<1>'
        while os.path.exists(newpath):
            repl += '_'
            newpath = re.sub(r'(\S*)(\s*)\.pptx', repl+'.pptx', path)
        os.rename(path, newpath)
    return newpath

remove_pptx_blank('/Users/hui/Desktop/2020-02-03')

int(2011)

ppt = Presentation('/Users/hui/Desktop/VEGF/RIC50 SELEX QC.pptx')
for slide in ppt.slides:
    temp = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            temp.append((shape.text, shape.top))
    if not temp:
        temp.append(("No Title", 0))
    temp.sort(key=lambda x: x[1])
    title = temp[0][0]
    match = date.search(title)
    if match:
        y,m,d,author = match.groups()
    else:
        y,m,d,author = '2011',1,1,None
    y = '20'+y if len(y)==2 else y
    time = datetime(int(y),int(m),int(d))
    print(time,author)

def printreplace(paths):
    for path in paths:
         newpath = re.sub(r'(\S*)(\s*)\.pptx', r'\g<1>.pptx', path)
         if path!=newpath:
             print(newpath)


datetime(20,1,1,)


uniTag = re.compile('<(?P<pt>tag|note|flag)>(?P<content>.*)</(?P=pt)>')
uniTag.search(notes)

a=uniTag.findall(notes)
dict(a)
a

uniTag.findall("<flag></flag>")

uniTag.findall('s')

list(filter(lambda x:x[0]=='tag',a))
a
tag = re.compile('<tag>(.*)</tag>')
note = re.compile('<note>(.*)</note>')
flag = re.compile('(?P<pt><flag>)(?P<content>.*)(?P<ft></flag>)')
notes = '<tag>this is a tag,</tag>\n<note>this is a note,</note><a>aga</a>'
re.sub(regx,r'<\g<pt>>'+'new'+r'</\g<pt>>',notes)
regx.search(notes).groups()


regx.search(notes)
flag.search("<flag> </flag>").group('content')

n = tag.search(notes)
n
n and n.group(1)
print(tag.search("<tag></tag>").group(1))
tag.search("<tag></tag>").group(1)
" a".strip()

flag.search(notes)
tag.search(notes)
note.search(notes).group(1)

notes.replace(note,'new')


ppt=Presentation("/Users/hui/Desktop/Apt Quant.pptx")
ppt.core_properties.revision
s1 = ppt.slides[0]
s2 = ppt.slides[3]
s1.has_notes_slide
s2.has_notes_slide
s1.notes_slide


text_frame = s1.notes_slide.notes_text_frame
text_frame.text
s2.notes_slide.notes_text_frame.text = '<tag>this is a tag,</tag>\n<note>this is a note,</note>'
s2.notes_slide.notes_text_frame.text

ppt.save("/Users/hui/Desktop/Apt Quant.pptx")
