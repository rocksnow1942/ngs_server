import hashlib
import sys
from pathlib import PurePath
import os,time
from datetime import datetime,timedelta
# filepath = os.path.dirname(__file__)
filepath = PurePath(__file__).parent.parent.parent
sys.path.append(str(filepath))
from dateutil import parser
from app import db
from app import create_app
from pptx import Presentation
from watchdog.observers import Observer
from watchdog.events import PatternMatchingEventHandler
from app.models import Slide, PPT, Project
from flask import current_app
import re
#TODO
#when trashing project, only save those slides with notes and tags.
#after a create is done, cleanup projects with no slides, cleanup powerpoints that doesn't have project name.
#clean up slides in that powerpoint that doesn't have tag or notes.
# the left over slides will be displayed in "deleted tab"


# this POSSIBLE_TAGS defined the possible tag fields writting to notes slide in pptx. 
POSSIBLE_TAGS = 'tag|note'

class PPT_Indexer():
    """
    module to sync pptx with database index. 
    """
    def __init__(self, log_file=""):
        app = create_app(keeplog=False)
        app.app_context().push()
        source_folder = app.config['PPT_SOURCE_FOLDER'] # this is the cloud station folder on the server. 
        target_folder = app.config['PPT_TARGET_FOLDER'] # this is the storage folder for ppt snap shots. 
        self.app=app
        self.source_folder = source_folder
        self.target_folder = target_folder
        self.log_file = log_file

    def write_log(self,content):
        """
        write a log
        """
        with open(self.log_file,'a') as f:
            f.write(f"{self.time} - " +content+'\n')

    def mirror_path(self, path):
        """
        create the mirror path for ppt inside source folder by remove folder structure between project folder and ppt file itself.
        assume pptx file path =  VEGF/FirstSELEX/First Try/SELEXppt.pptx; 
        mirror path would be: VEGF/SELEXppt 
        the mirror path will be the folder inside ppt snapshot folder. 
        """
        sf = PurePath(self.source_folder)
        # tf = PurePath(self.target_folder)
        rf = PurePath(path).relative_to(sf)
        return str(rf.parent/rf.stem).strip()

    def get_project_name(self,filepath):
        """get the first layer of path, that is the project folder name."""
        return PurePath(filepath).relative_to(self.source_folder).parts[0]

    def get_PPT_name(self,filepath):
        """get the ppt file name. remove space; because ppt snapshot folder will not have space. (this is caused by windows pptx making snapshot removing space.)"""
        return PurePath(filepath).stem.strip()

    @staticmethod
    def get_md5(file):
        "get the md5 of a pptx file."
        with open(file, 'rb') as afile:
            hasher = hashlib.md5()
            buf = afile.read()
            hasher.update(buf)
        return hasher.hexdigest()

    @staticmethod
    def get_revision(file):
        "get revision number of a pptx file"
        ppt = Presentation(file)
        return ppt.core_properties.revision
    @property
    def time(self):
        "format time"
        return datetime.now().strftime('%y/%m/%d %H:%M:%S')

    def create(self, file):
        """
        entry point for creating pptx log. 
        try to create the file by its path; this happens when regular editing of file.
        then by its md5; this happens when moving the same file to a different folder.
        then by create new file; this is when a brand new file is created. 
        write a log.
        """
        with self.app.app_context():
            try:
                if self.create_by_path(file):
                    self.write_log(f" Update by path {file}")
                elif self.create_by_md5((file)):
                    self.write_log(f" Update by md5 {file}")
                else:
                    self.create_new(file)
            except Exception as e:
                self.write_log(f" Create error:{file} - {e}")

    def create_by_path(self,file):
        """
        update the ppt log by file path. 
        first find ppt inside database using its path; then update its associated project and md5 and revision number; 
        then sync the slides.  
        """
        path = self.mirror_path(file)
        ppt = PPT.query.filter_by(path=path).first()
        if ppt: #and (ppt.project_id == None)
            project = self.create_project(file)
            ppt.project = project
            ppt.md5=self.get_md5(file)
            ppt.revision = self.get_revision(file)

            ppt=self.sync_slides(ppt,file)

            ppt.date = datetime.now()
            db.session.commit()
            return True
        return False

    def create_by_md5(self, file):
        "update by mod5, by only chaning to current path and project and ppt name."
        md5 = self.get_md5(file)
        ppt = PPT.query.filter_by(md5=md5).first()
        if ppt and (ppt.project_id == None):
            project = self.create_project(file)
            ppt.project = project
            ppt.name = self.get_PPT_name(file)
            ppt.path = self.mirror_path(file)
            ppt.date = datetime.now()
            db.session.commit()
            return True
        return False

    def create_new(self, file):
        """
        create a new PPT record. 
        """
        name = self.get_PPT_name(file)
        path = self.mirror_path(file)
        md5 = self.get_md5(file)
        project = self.create_project(file)
        revision = self.get_revision(file)
        try:
            ppt = PPT(name=name, path=path, project_id=project.id,
                      md5=md5, revision=revision)
            db.session.add(ppt)
            db.session.commit()
            self.sync_slides(ppt, file)
            db.session.commit()
            self.write_log(f" Create new {file}")
        except Exception as e:
            self.write_log(f" Create New {file} Error:{e}")
            raise e

    def parse_ppt(self,file):
        """
        use the presentation module to parse ppt slides into a list of dict. 
        [{page: , title: , date: , author: , body: , tag: , note: , (if tag and note exist) }]
        """
        ppt=Presentation(file)
        slides = []
        uniTags = re.compile(f'<(?P<pt>{POSSIBLE_TAGS})>(?P<content>.*)</(?P=pt)>') # parse tags in slide notes. like <tag></tag>

        # regx capture for date authoer: match start, match n space, match year/month/day group with non-word inbetween, match optional author with preceding blank. 
        dateAuthor = re.compile(r"\A\s*(?P<y>20\d{2}|\d{2})\W*(?P<m>[0]\d|1[0-2]|[1-9])\W*(?P<d>[0-2]\d|3[01]|\d)(?:\s*\((?P<author>[a-zA-Z\s]+)\))?")
        for page, slide in enumerate(ppt.slides):
            temp = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    temp.append((shape.text, shape.top))
            if not temp:
                temp.append(("No Title", 0))
            temp.sort(key=lambda x: x[1])
            title = temp[0][0]

            textNote = slide.notes_slide.notes_text_frame.text
            tagsFound = dict(uniTags.findall(textNote))
            matchDate = dateAuthor.search(title)
            if matchDate:
                y,m,d,author = matchDate.groups()
                y = '20'+y if len(y)==2 else y
                author = author.strip().upper() if author else author 
               
                try:
                    date = datetime(int(y),int(m),int(d)) 
                    now = datetime.now()
                    if (date> (now + timedelta(365 * 3)) ) or (date < datetime(2000,1,1)):
                        date = datetime.now()
                except:
                    date = slides[-1]['date'] if slides else datetime(
                        2011, 1, 1, 1, 1)
            else:
                parseresult = self.parse_date(title)
                if parseresult:
                    date = parseresult[0]
                    author = parseresult[1] or None
                else:
                    author = slides[-1]['author'] if slides else None
                    date = slides[-1]['date'] if slides else datetime(2011,1,1,1,1)
         
            tagsFound.update(page=page+1,
                               title=title, date=date,author=author,
                               body="\n".join(i[0] for i in temp[1:]))
            slides.append(tagsFound)
        return slides

    def parse_date(self,title):
        """
        parse the date by trial and error.
        not in use anymore.
        """
        date = ''
        for l in title:
            if l.isalpha():
                break
            else:
                date += l
        date=date.strip()
        for i in range(len(date),3,-1):
            totry = date[:i]
            try:
                _date = parser.parse(totry,ignoretz=True) # to avoid timezone error
                authorstring = title.replace(totry,"").strip()
                authorpattern = re.compile(
                    r"(?:^\((?P<author>[a-zA-Z\s]+)\))")
                match = authorpattern.match(authorstring) 
                if match: 
                    return _date, match.groups()[0].strip()
                else:
                    return _date, None
            except:
                continue
        return None

    def sync_slides(self,ppt,file):
        """
        make slide indexes to database logs. 
        if an slide with the same title and body already exist
        """
        slides = self.parse_ppt(file)
        oldslides=ppt.slides # get all the slides in the current ppt. 
        cur_slides = []

        # looping over the slides list from parsed list. 
        for s in slides:
            # find collection of slides that have the same title and body, and order them by page. 
            slide_collection = Slide.query.filter_by(ppt_id=ppt.id,title=s['title'],body=s['body']).order_by(Slide.page).all()
            found = False
            
            if slide_collection:
                # using collection to deal with the rare case where user create multiple new replicate slides at a random page. 
                # if that happens, will update the pages and note/tag in the order of their page. 
                # for the last page of the replicates, because it is not in the cur_slides collection, will be created as new slide. 
                for slide in slide_collection:
                    # if a slide in the collection is not been seen yet, update it accordingly. 
                    # the net effect is basically reindexing all slides with same title and body. 
                    if slide not in cur_slides:
                        slide.ppt=ppt
                        slide.page = s['page']
                        slide.date=s['date']
                        # update the tag and notes:
                        for t in POSSIBLE_TAGS.split('|'):
                            setattr(slide, t, s.get(t, None))
                        found = True
                        break
            if not found: 
                 # create new if not found. 
                slide = Slide(ppt_id=ppt.id, **s)
                db.session.add(slide)
            # maintain a log of slides that are already been indexed. 
            cur_slides.append(slide)

        # remove deleted slides; if the slide have note or tag, trash it by remove its ppt_id. 
        for slide in (set(oldslides)-set(cur_slides)): 
            if slide.note or slide.tag:
                slide.ppt_id=None
            else:
                db.session.delete(slide)
        db.session.commit()
        return ppt

    def create_project(self,file):
        """
        return corresponding project database entry given the file. 
        create a new project record if the project doesn't exist.
        """
        projectname = self.get_project_name(file)
        project = Project.query.filter_by(name=projectname).first()
        if not project:
            project=Project(name=projectname)
            db.session.add(project)
        project.date=datetime.now()
        db.session.commit()
        return project

    def delete(self, file):
        """
        entry point for deleting ppt. 
        this is always done by its file path. 
        when deleting, only remove the associated project_id, not actually deleting all the records. 
        """
        with self.app.app_context():
            filepath = self.mirror_path(file)
            ppt = PPT.query.filter_by(path=filepath).first()
            if ppt:
                ppt.project_id=None
                db.session.commit()
                self.write_log(
                    f" Delete PPT {file}")

    def move(self, src, dst):
        """
        this event rarely happens on Ubuntu and is not dealt with. 
        """
        self.write_log(f" Move {src} => {dst}")



class PPTX_Handler(PatternMatchingEventHandler):
    """
    watchdog event listner. 
    """
    def __init__(self, logger):
        super().__init__(patterns=["*.pptx", ], ignore_patterns=["*~$*", "*Conflict*"],
                         ignore_directories=False, case_sensitive=True)
        self.logger = logger

    def on_created(self, event):
        """
        on create event, start indexing the pptx. first remove trailing blank. 
        """
        t1=time.time()
        try:
            # first remove the bland before start indexing. 
            removedblank = remove_pptx_blank(event.src_path)
            self.logger.create(removedblank)
            t2=time.time()
            print("Create {} Done in {:.1f}".format(removedblank, t2-t1))
        except Exception as e:
            print(f" Create Error {removedblank}:{e}")

    def on_deleted(self, event):
        t1 = time.time()
        print("{} Delete {}".format(self.logger.time,event.src_path))
        try:
            self.logger.delete(event.src_path)
            t2 = time.time()
            print("Delete {} Done in {:.1f}".format(event.src_path, t2-t1))
        except Exception as e:
            self.logger.write_log(f"Deletion error:{e}")
            print(f"Delete Error {event.src_path}:{e}")

    def on_modified(self, event):
        print(f"Modify {event.src_path}")

    def on_moved(self, event):
        print(f"Rename {event.src_path} to {event.dest_path}")

def remove_pptx_blank(path):
    """
    remove trailing blank in pptx file path.
    """
    newpath = re.sub(r'(\S*)(\s*)\.pptx', r'\g<1>.pptx', path)
    if path != newpath and os.path.exists(path):
        repl = r'\g<1>'
        while os.path.exists(newpath):
            repl += '_'
            newpath = re.sub(r'(\S*)(\s*)\.pptx', repl+'.pptx', path)
        os.rename(path, newpath)
    return newpath
    

def StartWatch(source_folder, log_file):
    """
    start the observer. index every 10 seconds.
    """
    my_observer = Observer()
    logger = PPT_Indexer(log_file=log_file)

    my_observer.schedule(PPTX_Handler(logger=logger),
                         source_folder, recursive=True)
    my_observer.start()
    try:
        while True:
            time.sleep(10)
    except KeyboardInterrupt:
        my_observer.stop()
        my_observer.join()


def glob_pptx(path):
    """
    grab all the pptx files inside a folder path, ignoring all conflicts. 
    """
    result = []
    for root, dirs, files in os.walk(path):
        for file in files:
            if file.endswith(".pptx") and ('conflict' not in file.lower()) and (not file.startswith('~$')):
                  result.append(os.path.join(root, file))
    return result

def index_path(path,log_file):
    """
    rebuild index of each pptx file inside a path. 
    """
    logger = PPT_Indexer(log_file=log_file)
    pptxs=glob_pptx(path)
    for i in pptxs:
        logger.delete(i)
        logger.create(i)

def index_file(path):
    """
    index a specific file
    """
    app = create_app(keeplog=False)
    app.app_context().push()
    log_file = app.config['PPT_LOG_FILE']
    logger = PPT_Indexer(log_file=log_file)
    logger.delete(path)
    logger.create(path)

def reindex():
    """
    rebuild index for the current files. 
    """
    log_file = current_app.config['PPT_LOG_FILE']
    logger = PPT_Indexer(log_file=log_file)
    source_folder = PurePath(current_app.config['PPT_SOURCE_FOLDER'])
    paths = [str((source_folder)/PurePath(i.path))+".pptx" for i in PPT.query.all()]
    messages=[]
    for path in paths:
        messages.append(f'Sync <{path}>')
        if os.path.exists(path):
            try:
                logger.create(path)
                messages.append(f'Redindexed <{path}>')
            except Exception as e:
                 messages.append(f'Redindex Error <{path}> - <{e}>')
        else:
            try:
                logger.delete(path)
                messages.append(f'Deleted <{path}>')
            except Exception as e:
                 messages.append(f'Deletion Error <{path}> - <{e}>')
    return messages



if __name__ == "__main__":
    app = create_app(keeplog=False)
    app.app_context().push()
    source_folder = app.config['PPT_SOURCE_FOLDER']
    log_file = app.config['PPT_LOG_FILE']
    print(f'start watching {source_folder}')
    # index_path(
    #     '/Users/hui/Cloudstation/R&D/Projects/VEGF/VEGF Aptamer Characterization',log_file=log_file)

    StartWatch(source_folder, log_file)

# Create / Users/hui/Cloudstation/R &D/Projects/VEGF/Stemloop PD copy.pptx
# Delete / Users/hui/Cloudstation/R & D/Projects/VEGF/Stemloop PD copy.pptx
# Create / Users/hui/Cloudstation/R & D/Projects/VEGF/test slides.pptx
# Delete / Users/hui/Cloudstation/R & D/Projects/VEGF/test slides.pptx
# Create / Users/hui/Cloudstation/R & D/Projects/VEGF/test slides for ppt.pptx
# Delete / Users/hui/Cloudstation/R & D/Projects/VEGF/test slides for ppt.pptx
# Create / Users/hui/Cloudstation/R & D/Projects/VEGF/test slides for ppt.pptx
# Delete / Users/hui/Cloudstation/R & D/Projects/VEGF/test slides for ppt.pptx
# Create / Users/hui/Cloudstation/R & D/Projects/VEGF/test slides for ppt.pptx
# Delete / Users/hui/Cloudstation/R & D/Projects/VEGF/test slides for ppt.pptx
# Create / Users/hui/Cloudstation/R & D/Projects/VEGF/Test folder/test slides for ppt.pptx
